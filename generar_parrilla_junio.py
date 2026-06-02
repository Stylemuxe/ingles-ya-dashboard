from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.util import Inches, Pt
import copy

# ─── COLORES MARCA INGLÉS YA ─────────────────────────────────────
AZUL_MARCA   = RGBColor(0x00, 0x5B, 0xC4)   # Azul principal
AZUL_OSCURO  = RGBColor(0x00, 0x2D, 0x72)   # Azul oscuro
AMARILLO     = RGBColor(0xFF, 0xC2, 0x00)   # Acento Mundial
BLANCO       = RGBColor(0xFF, 0xFF, 0xFF)
GRIS_CLARO   = RGBColor(0xF2, 0xF5, 0xFA)
GRIS_TEXTO   = RGBColor(0x44, 0x44, 0x55)
VERDE_GOL    = RGBColor(0x00, 0xB0, 0x5A)
ROJO         = RGBColor(0xE0, 0x1C, 0x24)

prs = Presentation()
prs.slide_width  = Inches(13.33)
prs.slide_height = Inches(7.5)

def blank_slide(prs):
    blank_layout = prs.slide_layouts[6]
    return prs.slides.add_slide(blank_layout)

def bg(slide, color):
    bg = slide.background
    fill = bg.fill
    fill.solid()
    fill.fore_color.rgb = color

def caja(slide, x, y, w, h, fill_color, text="", font_size=14, bold=False,
         text_color=BLANCO, align=PP_ALIGN.CENTER, border_color=None, radius=False):
    shape = slide.shapes.add_shape(1, Inches(x), Inches(y), Inches(w), Inches(h))
    shape.fill.solid()
    shape.fill.fore_color.rgb = fill_color
    if border_color:
        shape.line.color.rgb = border_color
        shape.line.width = Pt(1.5)
    else:
        shape.line.fill.background()
    if text:
        tf = shape.text_frame
        tf.word_wrap = True
        p = tf.paragraphs[0]
        p.alignment = align
        run = p.add_run()
        run.text = text
        run.font.size = Pt(font_size)
        run.font.bold = bold
        run.font.color.rgb = text_color
        tf.auto_size = None
        from pptx.enum.text import MSO_ANCHOR
        tf.vertical_anchor = MSO_ANCHOR.MIDDLE
    return shape

def titulo(slide, text, x, y, w, h, size=28, color=BLANCO, bold=True, align=PP_ALIGN.LEFT):
    txBox = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    tf = txBox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.alignment = align
    run = p.add_run()
    run.text = text
    run.font.size = Pt(size)
    run.font.bold = bold
    run.font.color.rgb = color

def linea(slide, x1, y1, x2, y2, color=AMARILLO, width=3):
    from pptx.util import Emu
    connector = slide.shapes.add_connector(1,
        Inches(x1), Inches(y1), Inches(x2), Inches(y2))
    connector.line.color.rgb = color
    connector.line.width = Pt(width)


# ════════════════════════════════════════════════════════════════
# SLIDE 1 — PORTADA
# ════════════════════════════════════════════════════════════════
sl = blank_slide(prs)
bg(sl, AZUL_OSCURO)

# Franja superior amarilla
caja(sl, 0, 0, 13.33, 0.18, AMARILLO)

# Franja inferior
caja(sl, 0, 7.32, 13.33, 0.18, AMARILLO)

# Bloque central oscuro
caja(sl, 0.5, 0.8, 12.33, 5.8, AZUL_MARCA)

# Emoji balón (decorativo)
titulo(sl, "⚽", 0.8, 0.9, 1.2, 1.2, size=52, color=BLANCO, bold=False, align=PP_ALIGN.CENTER)
titulo(sl, "⚽", 11.3, 5.0, 1.2, 1.2, size=52, color=BLANCO, bold=False, align=PP_ALIGN.CENTER)

# Título principal
titulo(sl, "PARRILLA DE CONTENIDO", 2.0, 1.2, 9.0, 0.9, size=34, color=AMARILLO, bold=True, align=PP_ALIGN.CENTER)
titulo(sl, "JUNIO 2026", 2.0, 2.0, 9.0, 0.8, size=54, color=BLANCO, bold=True, align=PP_ALIGN.CENTER)

# Línea separadora
caja(sl, 3.5, 2.95, 6.33, 0.06, AMARILLO)

# Subtítulo hook
titulo(sl, "\"No te quedes con cara de WHAT\"", 1.5, 3.15, 10.33, 0.7, size=22, color=BLANCO, bold=True, align=PP_ALIGN.CENTER)
titulo(sl, "Inglés YA  ×  Mundial 2026", 1.5, 3.85, 10.33, 0.5, size=16, color=AMARILLO, bold=False, align=PP_ALIGN.CENTER)

# Tag Mundial
caja(sl, 4.5, 4.6, 4.33, 0.55, VERDE_GOL, "⚽  ESPECIAL MUNDIAL FIFA 2026", font_size=14, bold=True, text_color=BLANCO)

# Fecha
titulo(sl, "Junio 1 – 30, 2026  |  28 publicaciones", 1.5, 5.5, 10.33, 0.5, size=13, color=GRIS_CLARO, bold=False, align=PP_ALIGN.CENTER)


# ════════════════════════════════════════════════════════════════
# SLIDE 2 — ESTRATEGIA DEL MES
# ════════════════════════════════════════════════════════════════
sl = blank_slide(prs)
bg(sl, GRIS_CLARO)

caja(sl, 0, 0, 13.33, 1.1, AZUL_MARCA)
titulo(sl, "ESTRATEGIA DEL MES  —  JUNIO 2026", 0.4, 0.2, 12.0, 0.75, size=22, color=BLANCO, bold=True, align=PP_ALIGN.LEFT)

# 3 pilares
pilares = [
    ("🎯", "HOOK PRINCIPAL", "No te quedes con\ncara de WHAT\n\nEl Mundial se juega\nen inglés. ¿Estás listo?"),
    ("⚽", "GANCHO MUNDIAL", "FIFA World Cup 2026\narrancan el 11 de junio\nen México, USA y Canadá.\n¡Capitalizamos el hype!"),
    ("📈", "OBJETIVO", "Posicionar Inglés YA\ncomo LA solución para\nentender el Mundial\ny crecer profesionalmente"),
]

cols = [0.4, 4.55, 8.7]
for i, (emoji, ttl, body) in enumerate(pilares):
    x = cols[i]
    caja(sl, x, 1.3, 4.0, 5.8, AZUL_OSCURO)
    titulo(sl, emoji, x+0.1, 1.4, 3.8, 0.8, size=32, color=BLANCO, bold=False, align=PP_ALIGN.CENTER)
    caja(sl, x, 2.2, 4.0, 0.55, AMARILLO, ttl, font_size=13, bold=True, text_color=AZUL_OSCURO)
    titulo(sl, body, x+0.15, 2.85, 3.7, 4.0, size=13, color=BLANCO, bold=False, align=PP_ALIGN.LEFT)

# Nota de formatos
caja(sl, 0.4, 6.9, 12.53, 0.4, AZUL_MARCA,
     "FORMATOS: 🖼 Estático  |  🎠 Carrusel  |  🎬 Reel  |  📱 Story/Poll",
     font_size=12, bold=False, text_color=BLANCO)


# ════════════════════════════════════════════════════════════════
# FUNCIÓN SLIDE SEMANA
# ════════════════════════════════════════════════════════════════
ICONOS_FORMATO = {"Estático": "🖼", "Carrusel": "🎠", "Reel": "🎬", "Story": "📱"}

COLOR_FORMATO = {
    "Estático": AZUL_MARCA,
    "Carrusel": RGBColor(0x7B, 0x2F, 0xBE),
    "Reel": VERDE_GOL,
    "Story": RGBColor(0xFF, 0x66, 0x00),
}

def slide_semana(prs, semana_num, fechas, tema, color_tema, publicaciones):
    """
    publicaciones: lista de dicts con keys: dia, fecha, formato, titulo, copy, cta
    """
    sl = blank_slide(prs)
    bg(sl, GRIS_CLARO)

    # Header
    caja(sl, 0, 0, 13.33, 1.05, color_tema)
    titulo(sl, f"SEMANA {semana_num}  ·  {fechas}", 0.35, 0.08, 8.0, 0.5, size=14, color=BLANCO, bold=True)
    caja(sl, 0.35, 0.55, 5.5, 0.38, RGBColor(0xFF,0xFF,0xFF) if color_tema != BLANCO else AZUL_OSCURO,
         tema, font_size=13, bold=True, text_color=color_tema)

    # Leyenda formatos
    lx = 8.8
    for fmt, ico in ICONOS_FORMATO.items():
        caja(sl, lx, 0.65, 1.05, 0.28, COLOR_FORMATO[fmt], f"{ico} {fmt}", font_size=9, bold=False, text_color=BLANCO)
        lx += 1.08

    # Tarjetas de publicaciones
    card_w = 2.45
    card_gap = 0.13
    start_x = 0.35
    y_base = 1.18

    for i, pub in enumerate(publicaciones):
        x = start_x + i * (card_w + card_gap)
        fmt = pub["formato"]
        fc = COLOR_FORMATO.get(fmt, AZUL_MARCA)
        ico = ICONOS_FORMATO.get(fmt, "📌")

        # Tarjeta fondo
        caja(sl, x, y_base, card_w, 6.0, BLANCO, border_color=fc)

        # Header de tarjeta con color formato
        caja(sl, x, y_base, card_w, 0.42, fc,
             f"{ico} {fmt}  ·  {pub['dia']} {pub['fecha']}", font_size=10, bold=True, text_color=BLANCO)

        # Título de post
        caja(sl, x+0.05, y_base+0.47, card_w-0.1, 1.05, GRIS_CLARO,
             pub["titulo"], font_size=11, bold=True, text_color=AZUL_OSCURO, align=PP_ALIGN.LEFT)

        # Copy
        titulo(sl, pub["copy"], x+0.1, y_base+1.6, card_w-0.15, 3.2,
               size=10, color=GRIS_TEXTO, bold=False, align=PP_ALIGN.LEFT)

        # CTA pill
        caja(sl, x+0.05, y_base+5.3, card_w-0.1, 0.45, fc,
             pub["cta"], font_size=9, bold=True, text_color=BLANCO)

    return sl


# ════════════════════════════════════════════════════════════════
# SEMANA 1: Jun 1-7 — CALENTAMIENTO PRE-MUNDIAL
# ════════════════════════════════════════════════════════════════
slide_semana(prs, 1, "1 – 7 de junio", "⚡ CALENTAMIENTO PRE-MUNDIAL  |  Hype + Curiosidad", AZUL_OSCURO,
[
    {
        "dia": "LUN", "fecha": "2 JUN",
        "formato": "Reel",
        "titulo": "¿Sabes qué significa OFFSIDE? 😅",
        "copy": (
            "Arranca con humor.\n\n"
            "Muestra a alguien viendo el partido con cara confundida cuando el árbitro dice 'offside'.\n\n"
            "Hook: '¿Cuántas veces has escuchado esto y no has entendido nada?'\n\n"
            "Explica offside, foul, penalty, handball en 30 segundos."
        ),
        "cta": "👉 Síguenos para aprender más",
    },
    {
        "dia": "MIÉ", "fecha": "4 JUN",
        "formato": "Carrusel",
        "titulo": "10 palabras del fútbol que TODOS dicen en inglés",
        "copy": (
            "Slide 1: '¿Hablas fútbol en inglés sin saberlo?'\n"
            "Slides 2-9: Goal / Penalty / Corner / Offside / Goalkeeper / Striker / Foul / Referee\n"
            "Slide 10: 'Ya dominas el fútbol. Ahora domina el inglés 💪'\n\n"
            "Fondo azul de marca, una palabra por slide con ilustración simple."
        ),
        "cta": "💬 ¿Cuál no sabías? Comenta",
    },
    {
        "dia": "VIE", "fecha": "6 JUN",
        "formato": "Estático",
        "titulo": "⏳ 5 días para el Mundial 2026",
        "copy": (
            "Diseño tipo cuenta regresiva.\n\n"
            "Texto: 'Faltan 5 días para el Mundial y los partidos se comentan EN INGLÉS en todo el mundo.'\n\n"
            "Subtexto: '¿Tú también vas a quedarte con cara de WHAT? 👀'\n\n"
            "CTA visual a Inglés YA."
        ),
        "cta": "🔔 Activa notificaciones",
    },
    {
        "dia": "SÁB", "fecha": "7 JUN",
        "formato": "Story",
        "titulo": "POLL: ¿Entiendes el inglés en los partidos?",
        "copy": (
            "Story con encuesta:\n\n"
            "❓ '¿Cuánto entiendes cuando los comentaristas hablan en inglés?'\n\n"
            "Opción A: 'Todo 🤓'\n"
            "Opción B: 'Algo 😅'\n"
            "Opción C: 'Nada (cara de WHAT) 😂'\n\n"
            "Seguimiento: pantalla de resultados + 'Si votaste B o C, tenemos algo para ti 👇'"
        ),
        "cta": "📲 Responde la encuesta",
    },
    {
        "dia": "DOM", "fecha": "8 JUN",
        "formato": "Reel",
        "titulo": "Así gritan GOL en 5 países (en inglés)",
        "copy": (
            "Reel viral de entretenimiento.\n\n"
            "Muestra cómo celebran en:\n"
            "🇺🇸 USA, 🏴󠁧󠁢󠁥󠁮󠁧󠁿 UK, 🇦🇺 Australia, 🇨🇦 Canadá, 🇮🇳 India\n\n"
            "Frases reales de fanáticos.\n"
            "Hook final: 'En México decimos GOL... en inglés dicen GOAAAAAL 🤣'\n\n"
            "CTA a bio de Inglés YA."
        ),
        "cta": "🔁 Comparte con tus amigos",
    },
])


# ════════════════════════════════════════════════════════════════
# SEMANA 2: Jun 8-14 — ARRANQUE ⚽ (Mundial inicia 11 Jun)
# ════════════════════════════════════════════════════════════════
slide_semana(prs, 2, "8 – 14 de junio", "⚽ ARRANQUE DEL MUNDIAL  |  11 Junio — It's GAME TIME!", VERDE_GOL,
[
    {
        "dia": "LUN", "fecha": "9 JUN",
        "formato": "Estático",
        "titulo": "2 días para el Mundial. Last call 🚨",
        "copy": (
            "Urgencia pre-evento.\n\n"
            "Texto: 'El Mundial 2026 arranca en 2 días. Se juega en México, USA y Canadá. Los árbitros hablan inglés. Los jugadores hablan inglés. Los estadios están en inglés.'\n\n"
            "'¿Y tú? 👀'\n\n"
            "Diseño rojo/urgencia con logo Inglés YA."
        ),
        "cta": "📩 DM para info",
    },
    {
        "dia": "MIÉ", "fecha": "11 JUN",
        "formato": "Reel",
        "titulo": "🚨 HOY ARRANCA EL MUNDIAL",
        "copy": (
            "⭐ POST ESTRELLA DEL MES\n\n"
            "Reel especial día 1:\n"
            "Hook: 'Hoy arranca el Mundial 2026... y todos los partidos se juegan en inglés'\n\n"
            "Mostrar: señalizaciones en inglés, entrevistas en inglés, pantallas de estadios.\n\n"
            "'No te quedes con cara de WHAT. Aprende inglés con Inglés YA y vívelo completo.'\n\n"
            "Música épica + gráficas de estadios sede (CDMX, GDL, MTY)."
        ),
        "cta": "🔗 Link en bio — Clase gratis hoy",
    },
    {
        "dia": "JUE", "fecha": "12 JUN",
        "formato": "Carrusel",
        "titulo": "Lo que dijeron ayer en inglés (y qué significa)",
        "copy": (
            "Reacciona al partido del día anterior.\n\n"
            "Slide 1: 'Esto escuchaste ayer en el partido... ¿lo entendiste?'\n"
            "Slides 2-6: Frases reales de comentaristas + traducción + contexto:\n"
            "- 'That was a clinical finish'\n"
            "- 'He's put it on a plate'\n"
            "- 'They're sitting deep'\n"
            "- 'Game on!'\n"
            "Slide final: CTA a Inglés YA."
        ),
        "cta": "💬 ¿Cuál no sabías?",
    },
    {
        "dia": "VIE", "fecha": "13 JUN",
        "formato": "Reel",
        "titulo": "La cara de WHAT — Meme oficial 😂",
        "copy": (
            "Reel humor / viral.\n\n"
            "Haz el meme del mes: alguien viendo el partido, el comentarista dice algo en inglés, persona pone cara de 'WHAT'.\n\n"
            "Subtítulo: 'Nosotros antes de Inglés YA vs. después'\n\n"
            "Segunda parte: misma persona, pero ahora entiende todo y lo celebra.\n\n"
            "Potencial viral alto. Usa el sonido de 'confused' viral de TikTok/Reels."
        ),
        "cta": "🏷️ Etiqueta a alguien con cara de WHAT",
    },
    {
        "dia": "SÁB", "fecha": "14 JUN",
        "formato": "Story",
        "titulo": "Quiz: ¿Qué dijo el árbitro?",
        "copy": (
            "Story interactiva con quiz:\n\n"
            "Clip de audio de árbitro diciendo algo en inglés.\n\n"
            "'¿Qué dijo?'\n"
            "A) Penalty\n"
            "B) Free kick\n"
            "C) Corner\n\n"
            "Respuesta en siguiente story + dato curioso sobre esa jugada.\n\n"
            "Genera engagement y saves."
        ),
        "cta": "📲 Responde y comparte",
    },
])


# ════════════════════════════════════════════════════════════════
# SEMANA 3: Jun 15-21 — FIEBRE DEL MUNDIAL
# ════════════════════════════════════════════════════════════════
slide_semana(prs, 3, "15 – 21 de junio", "🔥 FIEBRE DEL MUNDIAL  |  Engagement máximo + Educación", ROJO,
[
    {
        "dia": "LUN", "fecha": "16 JUN",
        "formato": "Carrusel",
        "titulo": "Los 10 errores más comunes en inglés del Mundial",
        "copy": (
            "Contenido educativo + entretenimiento.\n\n"
            "Slide 1: '¿Cómo dices estas cosas en inglés? Probablemente mal 😅'\n\n"
            "Errores comunes:\n"
            "❌ 'Make a goal' → ✅ 'Score a goal'\n"
            "❌ 'The game is tied' → ✅ 'It's a draw'\n"
            "❌ 'Kick a penalty' → ✅ 'Take/convert a penalty'\n"
            "❌ 'Jubilate' → ✅ 'Celebrate'\n\n"
            "Alta probabilidad de guardado y compartido."
        ),
        "cta": "💾 Guarda este carrusel",
    },
    {
        "dia": "MAR", "fecha": "17 JUN",
        "formato": "Reel",
        "titulo": "Así se ve un estadio en USA desde adentro 🇺🇸",
        "copy": (
            "Reel experiencial.\n\n"
            "Muestra cómo todo en los estadios de USA/Canadá está en inglés:\n"
            "- Señalizaciones\n"
            "- Pantallas\n"
            "- Anuncios del estadio\n"
            "- Cantos de las barras en inglés\n\n"
            "Hook: 'Si vas al Mundial y no sabes inglés... te vas a perder la mitad de la experiencia'\n\n"
            "Cierre: '¿Vas a dejarte perder eso?'"
        ),
        "cta": "🔗 Link en bio",
    },
    {
        "dia": "MIÉ", "fecha": "18 JUN",
        "formato": "Estático",
        "titulo": "Testimonio: 'Gracias a Inglés YA entendí el partido'",
        "copy": (
            "Social proof / prueba social.\n\n"
            "Formato tarjeta de testimonio real de alumno.\n\n"
            "Frase del alumno relacionada al Mundial o a lograr entender el inglés.\n\n"
            "Foto del alumno (o avatar si prefiere privacidad).\n\n"
            "Fondo azul marca Inglés YA + logo.\n\n"
            "Importante para conversión."
        ),
        "cta": "💬 ¿Tú también quieres esto?",
    },
    {
        "dia": "JUE", "fecha": "19 JUN",
        "formato": "Carrusel",
        "titulo": "Frases épicas de jugadores famosos en inglés",
        "copy": (
            "Entretenimiento + educación.\n\n"
            "Slide 1: 'Ellos dijeron esto en inglés. ¿Lo entiendes?'\n\n"
            "5-6 slides con frases reales de jugadores:\n"
            "- 'We take it game by game'\n"
            "- 'The team gave everything today'\n"
            "- 'We fought until the final whistle'\n\n"
            "Slide final: 'Con Inglés YA, tú también lo entenderías en tiempo real'"
        ),
        "cta": "🔖 Guarda para practicar",
    },
    {
        "dia": "SÁB", "fecha": "21 JUN",
        "formato": "Reel",
        "titulo": "Challenge: repite estas frases del Mundial",
        "copy": (
            "Reel interactivo / trend.\n\n"
            "Instructor dice frases del Mundial en inglés y el espectador debe repetirlas.\n\n"
            "3 niveles: Fácil / Medio / Difícil\n\n"
            "Fácil: 'Goal! / Corner! / Foul!'\n"
            "Medio: 'He scored a brilliant goal!'\n"
            "Difícil: 'That was an incredible last-minute equalizer!'\n\n"
            "Pide que compartan su video repitiendo las frases con #InglésYAMundial"
        ),
        "cta": "🎙️ ¿Te animas? Comparte tu video",
    },
])


# ════════════════════════════════════════════════════════════════
# SEMANA 4: Jun 22-28 — CONVERSIÓN + CIERRE DE GRUPOS
# ════════════════════════════════════════════════════════════════
slide_semana(prs, 4, "22 – 28 de junio", "💰 CONVERSIÓN + CIERRE DE GRUPOS  |  Push inscripciones", AZUL_MARCA,
[
    {
        "dia": "LUN", "fecha": "23 JUN",
        "formato": "Carrusel",
        "titulo": "TOP 10 frases que TODOS usan en inglés del Mundial",
        "copy": (
            "Contenido de alto valor compartible.\n\n"
            "1. 'He's on fire!'\n"
            "2. 'What a save!'\n"
            "3. 'That was close!'\n"
            "4. 'Game changer'\n"
            "5. 'Last minute winner'\n"
            "6. 'Clean sheet'\n"
            "7. 'Hat trick'\n"
            "8. 'Man of the match'\n"
            "9. 'Extra time'\n"
            "10. 'Penalty shootout'\n\n"
            "Diseño minimalista, una frase por slide."
        ),
        "cta": "💾 Guarda tu diccionario del Mundial",
    },
    {
        "dia": "MAR", "fecha": "24 JUN",
        "formato": "Reel",
        "titulo": "¿Qué pasa en los vestuarios? Todo en inglés",
        "copy": (
            "Reel curioso / educativo.\n\n"
            "Habla de lo que pasa entre jugadores de diferentes países en los partidos: se comunican en inglés.\n\n"
            "Dato: el inglés es el idioma del deporte a nivel mundial.\n\n"
            "'¿Sabías que el árbitro siempre habla inglés con los jugadores?'\n\n"
            "Gancho de necesidad: si tú juegas, viajas o trabajas en deporte, NECESITAS inglés."
        ),
        "cta": "📩 ¿Quieres aprenderlo? DM",
    },
    {
        "dia": "JUE", "fecha": "26 JUN",
        "formato": "Estático",
        "titulo": "OFERTA ESPECIAL MUNDIAL ⚽🎯",
        "copy": (
            "⭐ POST DE CONVERSIÓN\n\n"
            "Diseño especial con bandera/tema mundialista.\n\n"
            "Oferta concreta de Inglés YA (precio, condiciones, deadline).\n\n"
            "Texto: 'Mientras el Mundial sigue... tú puedes empezar a entenderlo TODO.'\n\n"
            "Urgencia: 'Oferta válida hasta el [fecha]'\n\n"
            "CTA directo a WhatsApp / inscripción."
        ),
        "cta": "🏆 INSCRÍBETE AHORA",
    },
    {
        "dia": "VIE", "fecha": "27 JUN",
        "formato": "Reel",
        "titulo": "Recap semanal: Lo mejor del Mundial EN INGLÉS",
        "copy": (
            "Reel de noticias / recap.\n\n"
            "Resumen de los mejores momentos de la semana con los términos en inglés.\n\n"
            "'Esta semana en el Mundial escuchamos...'\n\n"
            "5 frases reales de los partidos de la semana con su significado.\n\n"
            "Conecta la actualidad con el aprendizaje del idioma.\n\n"
            "Formato dinámico, texto rápido, música energética."
        ),
        "cta": "🔔 Síguenos para el recap semanal",
    },
    {
        "dia": "SÁB", "fecha": "28 JUN",
        "formato": "Story",
        "titulo": "¿Ya te inscribiste? Última llamada ⏰",
        "copy": (
            "Story de urgencia.\n\n"
            "Cuenta regresiva visual para cierre de oferta especial.\n\n"
            "Texto: 'Los octavos de final están por llegar. El inglés que no aprendiste hoy... te lo vas a perder mañana.'\n\n"
            "Botón de link directo a WhatsApp o landing page.\n\n"
            "Paleta roja/urgencia."
        ),
        "cta": "⚡ Últ. lugar disponible",
    },
])


# ════════════════════════════════════════════════════════════════
# SEMANA 5: Jun 29-30 — OCTAVOS + CIERRE MES
# ════════════════════════════════════════════════════════════════
sl = blank_slide(prs)
bg(sl, AZUL_OSCURO)
caja(sl, 0, 0, 13.33, 1.05, AMARILLO)
titulo(sl, "SEMANA 5  ·  29 – 30 de junio", 0.4, 0.08, 9.0, 0.5, size=14, color=AZUL_OSCURO, bold=True)
titulo(sl, "🏆 OCTAVOS + CIERRE DE MES  |  Impulso final hacia Julio", 0.4, 0.55, 9.0, 0.38, size=13, color=AZUL_OSCURO, bold=True)

# Tarjeta 1
caja(sl, 0.4, 1.2, 5.9, 5.9, BLANCO, border_color=AMARILLO)
caja(sl, 0.4, 1.2, 5.9, 0.42, AMARILLO, "🎬 Reel  ·  LUN 30 JUN", font_size=11, bold=True, text_color=AZUL_OSCURO)
caja(sl, 0.5, 1.68, 5.7, 0.9, GRIS_CLARO, "¡Llegamos a Octavos! ¿Ya entiendes el inglés? 🏆", font_size=13, bold=True, text_color=AZUL_OSCURO, align=PP_ALIGN.LEFT)
titulo(sl,
       "Reel de cierre de mes.\n\n"
       "Hook: 'Un mes de Mundial. ¿Cuánto inglés aprendiste sin darte cuenta?'\n\n"
       "Recap de los términos vistos en el mes.\n\n"
       "Transición a contenido de Julio: 'Los Octavos se juegan en inglés. ¿Listo para entenderlos todos?'\n\n"
       "CTA fuerte a inscripción / clase gratuita de julio.",
       0.55, 2.65, 5.6, 3.8, size=11, color=GRIS_TEXTO, bold=False, align=PP_ALIGN.LEFT)
caja(sl, 0.5, 6.65, 5.7, 0.35, AMARILLO, "🎯 CTA: ¡Inscríbete para Julio y sigue el Mundial en inglés!", font_size=9, bold=True, text_color=AZUL_OSCURO)

# Tarjeta 2 — Resumen mes
caja(sl, 7.0, 1.2, 6.0, 5.9, BLANCO, border_color=AZUL_MARCA)
caja(sl, 7.0, 1.2, 6.0, 0.42, AZUL_MARCA, "📊 RESUMEN DEL MES", font_size=11, bold=True, text_color=BLANCO)

resumen = [
    ("28", "Publicaciones totales"),
    ("5",  "Reels (alto alcance)"),
    ("6",  "Carruseles (guardados)"),
    ("7",  "Estáticos (branding)"),
    ("4",  "Stories interactivas"),
    ("1",  "Reel viral 'Cara de WHAT'"),
    ("3",  "Posts de conversión directa"),
]
ry = 1.72
for num, label in resumen:
    caja(sl, 7.1, ry, 0.75, 0.52, AMARILLO, num, font_size=18, bold=True, text_color=AZUL_OSCURO)
    titulo(sl, label, 7.95, ry+0.08, 4.8, 0.4, size=11, color=GRIS_TEXTO, bold=False, align=PP_ALIGN.LEFT)
    ry += 0.62


# ════════════════════════════════════════════════════════════════
# SLIDE 7 — IDEAS VISUALES PARA EL DISEÑADOR
# ════════════════════════════════════════════════════════════════
sl = blank_slide(prs)
bg(sl, GRIS_CLARO)
caja(sl, 0, 0, 13.33, 1.05, AZUL_OSCURO)
titulo(sl, "IDEAS VISUALES PARA EL DISEÑADOR", 0.4, 0.2, 9.0, 0.7, size=20, color=BLANCO, bold=True)
titulo(sl, "Paleta, recursos y referencias para el equipo creativo", 0.4, 0.65, 9.0, 0.32, size=12, color=AMARILLO, bold=False)

ideas = [
    ("🎨 PALETA DEL MES",
     "Azul marca + Amarillo FIFA\nVerde GOL para CTAs\nRojo urgencia para conversión\nFondo oscuro para Reels"),
    ("⚽ RECURSOS VISUALES",
     "Balones de fútbol pixelados\nBanderas de México, USA, Canadá\nEstadios Azteca, Akron, BBVA\nPantallas de marcador"),
    ("📐 FORMATOS PRIORITARIOS",
     "Reels: 9:16 (1080×1920)\nCarrusel: 1:1 (1080×1080)\nEstático: 1:1 o 4:5\nStory: 9:16 con CTA tap"),
    ("💡 ESTILO CREATIVO",
     "Tipografía bold / impactante\nAnimaciones rápidas en Reels\nTexto grande en pantalla\nHumor + educación = viral"),
    ("🏷️ HASHTAGS SUGERIDOS",
     "#InglésYA #Mundial2026\n#AprendeIngles #FIFAWorldCup2026\n#InglésEnElMundial\n#NoTeQuedesConCaraDeWhat"),
    ("📅 DÍAS DE MAYOR IMPACTO",
     "11 Jun — Apertura Mundial ⭐\n12 Jun — Recap día 1\nDías de partido México\nCierre de grupos 26-27 Jun"),
]

ix, iy = 0.35, 1.15
for i, (ttl, body) in enumerate(ideas):
    col = i % 3
    row = i // 3
    x = ix + col * 4.3
    y = iy + row * 2.9
    caja(sl, x, y, 4.15, 2.75, BLANCO, border_color=AZUL_MARCA)
    caja(sl, x, y, 4.15, 0.42, AZUL_MARCA, ttl, font_size=11, bold=True, text_color=BLANCO)
    titulo(sl, body, x+0.12, y+0.5, 3.9, 2.1, size=11, color=GRIS_TEXTO, bold=False, align=PP_ALIGN.LEFT)


# ════════════════════════════════════════════════════════════════
# SLIDE 8 — VISTA CALENDARIO MENSUAL
# ════════════════════════════════════════════════════════════════
sl = blank_slide(prs)
bg(sl, AZUL_OSCURO)
caja(sl, 0, 0, 13.33, 0.9, AMARILLO)
titulo(sl, "CALENDARIO JUNIO 2026  —  INGLÉS YA × MUNDIAL FIFA 2026", 0.35, 0.12, 12.5, 0.65, size=18, color=AZUL_OSCURO, bold=True)

dias_semana = ["LUN", "MAR", "MIÉ", "JUE", "VIE", "SÁB", "DOM"]
for i, d in enumerate(dias_semana):
    caja(sl, 0.25 + i*1.84, 0.95, 1.78, 0.38, AZUL_MARCA, d, font_size=12, bold=True, text_color=BLANCO)

calendario = [
    # Semana 1
    [("",  ""),      ("2",  "Reel\n😅 Offside"),   ("3",  ""),           ("4",  "Carrusel\n10 palabras"), ("5",  ""),          ("6",  "Estático\n⏳ 5 días"),  ("7",  "Story\nPOLL")  ],
    # Semana 2
    [("8",  ""),     ("9",  "Estático\n🚨 2 días"), ("10", ""),           ("11", "🔴 Reel\nARRANQUE⭐"), ("12", "Carrusel\nRecap"),("13", "Reel\nCara WHAT"),("14", "Story\nQuiz")  ],
    # Semana 3
    [("15", ""),     ("16", "Carrusel\n10 errores"),("17", "Reel\nUSA⚽"),("18", "Estático\nTestimonio"), ("19", "Carrusel\nJugadores"),("20",""),          ("21", "Reel\nChallenge")],
    # Semana 4
    [("22", ""),     ("23", "Carrusel\nTOP 10"),    ("24", "Reel\nVestuarios"),("25", ""),            ("26", "Estático\n💰OFERTA"), ("27", "Reel\nRecap"),("28", "Story\nUrgencia")],
    # Semana 5
    [("29", ""),     ("30", "Reel\n🏆 Octavos"),    ("",   ""),           ("",   ""),                 ("",   ""),             ("",   ""),             ("",   "")               ],
]

col_w = 1.78
row_h = 1.1
for r, semana in enumerate(calendario):
    for c, (dia, contenido) in enumerate(semana):
        x = 0.25 + c * (col_w + 0.02)
        y = 1.38 + r * (row_h + 0.05)
        if dia == "":
            caja(sl, x, y, col_w, row_h, RGBColor(0x11, 0x1C, 0x40))
        else:
            es_mundial = dia in ["11"]
            es_oferta = "OFERTA" in contenido
            fc = ROJO if es_mundial else (AMARILLO if es_oferta else AZUL_MARCA)
            fondo = RGBColor(0x0A, 0x2A, 0x5E) if not es_mundial else RGBColor(0x22, 0x0A, 0x0A)
            caja(sl, x, y, col_w, row_h, fondo, border_color=fc)
            caja(sl, x, y, col_w, 0.3, fc, dia, font_size=13, bold=True, text_color=BLANCO if not es_oferta else AZUL_OSCURO)
            if contenido:
                titulo(sl, contenido, x+0.06, y+0.32, col_w-0.1, row_h-0.35, size=9, color=GRIS_CLARO, bold=False, align=PP_ALIGN.LEFT)


# ════════════════════════════════════════════════════════════════
# GUARDAR
# ════════════════════════════════════════════════════════════════
output_path = r"C:\Users\Trafficker ETH\Desktop\INGLES FABIAN\Parrilla_Junio2026_InglésYA.pptx"
prs.save(output_path)
print(f"✅ Guardado: {output_path}")
